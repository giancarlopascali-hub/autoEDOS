import os
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt

def generate_contours():
    x = np.linspace(-3, 3, 100)
    y = np.linspace(-3, 3, 100)
    X, Y = np.meshgrid(x, y)
    
    # Create two peaks (a local and a global)
    Z = 3 * (1 - X)**2 * np.exp(-(X**2) - (Y + 1)**2) - \
        10 * (X / 5 - X**3 - Y**5) * np.exp(-X**2 - Y**2) - \
        (1/3) * np.exp(-(X + 1)**2 - Y**2)
    return X, Y, Z

def create_ovat_vs_doe_plot():
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))
    X, Y, Z = generate_contours()
    
    # OVAT Plot
    ax1.contourf(X, Y, Z, levels=20, cmap='viridis', alpha=0.8)
    ax1.set_title("OVAT Strategy", fontsize=14, pad=10)
    
    # OVAT Path - Starts far away, moves X, then stops at local optimum, then Y
    ovat_x1 = np.linspace(-2.5, 0.5, 5)
    ovat_y1 = [-2.5] * 5
    ax1.plot(ovat_x1, ovat_y1, 'w-o', linewidth=3, markersize=8, label='Step 1 (Vary X)')
    
    ovat_x2 = [0.5] * 5
    ovat_y2 = np.linspace(-2.5, 1, 5)
    ax1.plot(ovat_x2, ovat_y2, 'r-o', linewidth=3, markersize=8, label='Step 2 (Vary Bottom)')
    
    ax1.scatter([0.5], [1], c='red', s=200, marker='*', zorder=10, label='Local Maximum')
    ax1.legend(loc='lower left')
    ax1.set_xlabel("Parameter 1")
    ax1.set_ylabel("Parameter 2")
    
    # DoE Plot
    ax2.contourf(X, Y, Z, levels=20, cmap='viridis', alpha=0.8)
    ax2.set_title("DoE (Factorial + Center)", fontsize=14, pad=10)
    
    # Generate points across the space
    doe_xs = [-2, -2, 2, 2, 0]
    doe_ys = [-2, 2, -2, 2, 0]
    ax2.scatter(doe_xs, doe_ys, c='white', s=100, edgecolors='k', zorder=5, label='Design Points')
    
    # Easily finds global maxima vicinity
    ax2.scatter([-0.5], [1.5], c='gold', s=300, marker='*', zorder=10, label='Global Maximum Discovered')
    ax2.legend(loc='lower left')
    ax2.set_xlabel("Parameter 1")
    ax2.set_ylabel("Parameter 2")
    
    plt.tight_layout()
    plt.savefig('ovat_vs_doe.png', dpi=150)
    plt.close()


def create_doe_metrics_plot():
    fig = plt.figure(figsize=(12, 5))
    
    # Left Side: Orthogonality
    ax1 = fig.add_subplot(121)
    ax1.set_title("Orthogonality (Correlation = 0)", fontsize=14, pad=10)
    ax1.set_xlim(0, 1.2)
    ax1.set_ylim(0, 1.2)
    
    # Orthogonal vectors
    ax1.quiver(0.1, 0.1, 0, 0.8, angles='xy', scale_units='xy', scale=1, color='green', width=0.015, label='Factor 1')
    ax1.quiver(0.1, 0.1, 0.8, 0, angles='xy', scale_units='xy', scale=1, color='blue', width=0.015, label='Factor 2')
    ax1.quiver(0.4, 0.4, 0.5, 0.5, angles='xy', scale_units='xy', scale=1, color='red', width=0.015, alpha=0.5, label='Correlated (Poor)')
    
    ax1.text(0.15, 0.4, "Independent\nEstimation", color='green', fontsize=12)
    ax1.text(0.6, 0.6, "Confounding", color='red', fontsize=12)
    ax1.grid(True, linestyle='--', alpha=0.5)
    ax1.legend()
    
    # Right Side: Curvature (3D)
    ax2 = fig.add_subplot(122, projection='3d')
    ax2.set_title("Curvature Detection", fontsize=14, pad=10)
    
    x = np.linspace(-1, 1, 20)
    y = np.linspace(-1, 1, 20)
    X, Y = np.meshgrid(x, y)
    Z = -(X**2 + Y**2)
    
    ax2.plot_surface(X, Y, Z, cmap='coolwarm', alpha=0.6)
    
    # Show linear plane entirely missing middle
    Z_lin = np.zeros_like(X) - 1
    ax2.plot_wireframe(X, Y, Z_lin, color='gray', alpha=0.4, label='Linear Model (No Center)')
    
    # Highlight Center Point
    ax2.scatter(0, 0, 0, c='black', s=100, marker='o', zorder=10)
    ax2.text(0, 0, 0.5, "Center Point required\nfor curve estimation", color='black', fontsize=11, ha='center')
    
    plt.tight_layout()
    plt.savefig('doe_metrics.png', dpi=150)
    plt.close()


def create_bo_concept_plot():
    np.random.seed(42)
    plt.figure(figsize=(10, 5))
    
    x = np.linspace(0, 10, 100)
    true_y = np.sin(x) + np.sin(2 * x)
    
    mean = true_y + 0.3 * np.sin(3*x)
    std = np.interp(x, [0, 2, 4, 6, 8, 10], [0.1, 1.2, 0.2, 0.9, 0.1, 1.0])
    
    # Acq func = mean + 1.96*std
    acq = mean + 1.96*std
    
    plt.plot(x, mean, 'b-', linewidth=3, label='Surrogate Mean Prediction')
    plt.fill_between(x, mean - 1.96*std, mean + 1.96*std, color='blue', alpha=0.2, label='Uncertainty (Std Dev)')
    
    plt.plot(x, acq, 'g--', linewidth=2, label='Acquisition Function (UCB)')
    
    # High Exploitation Area
    plt.annotate('Exploitation Target\n(High Mean)', xy=(2, 1.5), xytext=(0.5, 3), 
                 arrowprops=dict(facecolor='black', arrowstyle='->'), fontsize=12)
    
    # High Exploration Area
    plt.annotate('Exploration Target\n(High Uncertainty)', xy=(6.5, 1.5), xytext=(7, 3), 
                 arrowprops=dict(facecolor='black', arrowstyle='->'), fontsize=12)
    
    plt.title("Bayesian Gaussian Processes: Exploitation vs Exploration", fontsize=16)
    plt.legend(loc='lower right')
    plt.grid(True, linestyle='--', alpha=0.4)
    plt.tight_layout()
    plt.savefig('bo_concept.png', dpi=150)
    plt.close()


def create_edbo_vs_edos_plot():
    fig = plt.figure(figsize=(12, 5))
    
    # Left: EDBO+
    ax1 = fig.add_subplot(121, projection='3d')
    ax1.set_title("EDBO+ Categorical Approach\n(One-Hot Random Jumps)", fontsize=14, pad=10)
    
    x = np.linspace(-1, 1, 10)
    y = np.linspace(-1, 1, 10)
    X, Y = np.meshgrid(x, y)
    Z1 = np.zeros_like(X)
    Z2 = np.ones_like(X) * 2
    
    ax1.plot_surface(X, Y, Z1, color='lightblue', alpha=0.5)
    ax1.plot_surface(X, Y, Z2, color='lightgreen', alpha=0.5)
    
    ax1.plot([-0.5, 0.5], [0, 0], [0, 2], 'r-->', linewidth=2, label='Random Plane Jump')
    ax1.text(0, 0, 1, "No Gradient between\nindependent categories", color='red', fontsize=10)
    
    ax1.set_zticklabels(['Cat A', 'Cat B'])
    
    # Right: EDOS
    ax2 = fig.add_subplot(122)
    ax2.set_title("EDOS Discrete Mapping\n(Gradient Clamping)", fontsize=14, pad=10)
    
    x2 = np.linspace(0, 10, 100)
    y2 = -((x2 - 5)**2) + 10
    
    ax2.plot(x2, y2, 'b-', linewidth=2, label='Underlying Continuous Surrogate')
    
    # Grid lines
    for grid in [1, 3, 5, 7, 9]:
        ax2.axvline(x=grid, color='gray', linestyle='--', alpha=0.3)
    
    ax2.annotate('', xy=(4.5, 9.7), xytext=(2, 1), arrowprops=dict(facecolor='blue', width=2, headwidth=10))
    ax2.scatter([5], [10], color='green', s=150, zorder=5)
    
    # Snapping arrow
    ax2.annotate('Mathematical Snapping', xy=(5, 10), xytext=(7, 8), arrowprops=dict(facecolor='green', arrowstyle='->'))
    
    ax2.set_xlabel("Discrete Ordered Variable (e.g. integer amounts)")
    ax2.legend(loc='lower left')
    
    plt.tight_layout()
    plt.savefig('edbo_vs_edos.png', dpi=150)
    plt.close()


def create_sa_schema():
    fig, ax = plt.subplots(figsize=(10, 4))
    
    ax.axis('off')
    
    # Center Database
    rect = plt.Rectangle((0, 0.4), 2, 0.4, facecolor='lightblue', edgecolor='black', linewidth=2)
    ax.add_patch(rect)
    ax.text(1, 0.6, "100% Dataset", ha='center', va='center', fontsize=14, weight='bold')
    
    # Splits
    colors = ['orange', 'orange', 'orange', 'green', 'green']
    for i in range(5):
        y_pos = 0.8 - (i * 0.15)
        # Train
        ax.add_patch(plt.Rectangle((3, y_pos), 1.6, 0.1, facecolor='orange', alpha=0.7))
        # Test
        ax.add_patch(plt.Rectangle((4.6, y_pos), 0.4, 0.1, facecolor='green', alpha=0.7))
        
        ax.text(4, y_pos + 0.05, f"Split {i+1} (Train 80%)", ha='center', va='center', fontsize=10)
        
        # Arrow mapping
        ax.plot([5.1, 5.8], [y_pos + 0.05, y_pos + 0.05], 'k->', linewidth=1)
        ax.text(6.3, y_pos + 0.05, f"$R^2$ Score {i+1}", ha='center', va='center', fontsize=10)
        
    ax.plot([6.3, 8], [0.5, 0.5], 'k->', linewidth=3)
    
    ax.add_patch(plt.Rectangle((8.1, 0.35), 1.6, 0.3, facecolor='violet', edgecolor='black', linewidth=2))
    ax.text(8.9, 0.55, "Average Scores\n=\nReliability Metric", ha='center', va='center', fontsize=12, weight='bold')

    plt.tight_layout()
    plt.savefig('sa_bootstrapping.png', dpi=150)
    plt.close()

def add_full_image_slide(prs, title, image_path, text_body=""):
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.title
    title_box.text = title
    
    if text_body:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text_body
        p.font.size = Pt(18)
        img_top = Inches(2.2)
    else:
        img_top = Inches(1.5)
        
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.5), img_top, width=Inches(9))
    return slide

def generate_presentation():
    # 1. Generate new visual charts
    print("Generating Matplotlib Graphics...")
    create_ovat_vs_doe_plot()
    create_doe_metrics_plot()
    create_bo_concept_plot()
    create_edbo_vs_edos_plot()
    create_sa_schema()

    # 2. Build Presentation
    prs = Presentation()
    
    # Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Engineered Design and Optimization System (EDOS)"
    slide.placeholders[1].text = "Visual Tour: Design of Experiments & Bayesian Analytics"

    # OVAT vs DoE
    add_full_image_slide(prs, "The Pitfalls of One-Variable-At-A-Time (OVAT)", 
                         "ovat_vs_doe.png",
                         "Mathematical Reality: Optimizing one axis at a time misses interaction peaks and crashes prematurely into local optima. DoE samples geometrically to map the entire topography.")

    # DoE Metrics
    add_full_image_slide(prs, "Evaluating Mathematical Design Quality", 
                         "doe_metrics.png",
                         "• Orthogonality allows isolation of pure variable effects without confounding.\n• Curvature requires center-points to mathematically curve the regression plane.")
                         
    # BO Concept
    add_full_image_slide(prs, "Bayesian Optimization: Gaussian Processes", 
                         "bo_concept.png",
                         "Acquisition functions strictly balance sampling High-Predicted-Mean (Exploitation) versus High-Uncertainty areas (Exploration).")                         

    # EDBO vs EDOS
    add_full_image_slide(prs, "Navigating Non-Continuous Variables", 
                         "edbo_vs_edos.png",
                         "Instead of treating integers as disconnected categories (One-Hot), EDOS forces surrogate gradients to climb connected scales, locking to grid points immediately prior to evaluation.")                         

    # Benchmarks
    if os.path.exists("benchmark_comparison_plots.png"):
        add_full_image_slide(prs, "Empirical Benchmarking: EDOS Internal vs EDBO+", 
                             "benchmark_comparison_plots.png",
                             "Thanks to continuous-gradient locking, EDOS achieves smoother convergence trajectories towards the True Optima compared to random categorical rotations.")                         

    # SA Bootstrapping
    add_full_image_slide(prs, "Statistical Analysis (SA): Bootstrapping Engine", 
                         "sa_bootstrapping.png",
                         "To mathematically guard against Model Overfitting, the SA module utilizes rigorous 80/20 ShuffleSplit validation, punishing predictors that memorize data.")                         

    prs.save("EDOS_Presentation_Visual.pptx")
    print("Saved -> EDOS_Presentation_Visual.pptx")

if __name__ == "__main__":
    generate_presentation()
